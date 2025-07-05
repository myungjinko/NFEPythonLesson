#!/usr/bin/env python3
"""
NFE Python Course - Markdown to PowerPoint Converter

This script converts the Markdown presentation to PowerPoint format.
Requires: pip install python-pptx

Usage: python convert_to_powerpoint.py
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_powerpoint_from_markdown(markdown_file, output_file):
    """Convert Markdown presentation to PowerPoint"""
    
    # Read the markdown file
    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Create a new presentation
    prs = Presentation()
    
    # Split content into slides (separated by ---)
    slides_content = content.split('---')
    
    for slide_content in slides_content:
        slide_content = slide_content.strip()
        if not slide_content:
            continue
            
        # Create a new slide
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Extract title and content
        lines = slide_content.split('\n')
        title = ""
        content_lines = []
        
        for line in lines:
            if line.startswith('# ') and not title:
                title = line[2:].strip()
            elif line.startswith('## ') and not title:
                title = line[3:].strip()
            else:
                content_lines.append(line)
        
        # Set title
        if title:
            title_shape = slide.shapes.title
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set content
        content_text = '\n'.join(content_lines).strip()
        if content_text:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            # Process content line by line
            for line in content_lines:
                line = line.strip()
                if not line:
                    continue
                    
                # Handle different content types
                if line.startswith('### '):
                    # Subheading
                    p = text_frame.add_paragraph()
                    p.text = line[4:]
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(52, 73, 94)
                elif line.startswith('- ') or line.startswith('* '):
                    # Bullet point
                    p = text_frame.add_paragraph()
                    p.text = line[2:]
                    p.font.size = Pt(18)
                    p.level = 0
                elif line.startswith('  - ') or line.startswith('  * '):
                    # Sub-bullet point
                    p = text_frame.add_paragraph()
                    p.text = line[4:]
                    p.font.size = Pt(16)
                    p.level = 1
                elif line.startswith('```python'):
                    # Code block start
                    continue
                elif line.startswith('```'):
                    # Code block end
                    continue
                elif line.startswith('**') and line.endswith('**'):
                    # Bold text
                    p = text_frame.add_paragraph()
                    p.text = line[2:-2]
                    p.font.size = Pt(18)
                    p.font.bold = True
                elif line.startswith('1. ') or line.startswith('2. ') or line.startswith('3. ') or line.startswith('4. ') or line.startswith('5. '):
                    # Numbered list
                    p = text_frame.add_paragraph()
                    p.text = line[3:]
                    p.font.size = Pt(18)
                    p.level = 0
                else:
                    # Regular text
                    if line and not line.startswith('#'):
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.font.size = Pt(18)
    
    # Save the presentation
    prs.save(output_file)
    print(f"PowerPoint presentation saved as: {output_file}")

def main():
    """Main function"""
    input_file = "NFE_Python_Course_Presentation.md"
    output_file = "NFE_Python_Course_Presentation.pptx"
    
    try:
        create_powerpoint_from_markdown(input_file, output_file)
        print("✅ Conversion completed successfully!")
        print(f"📁 Output file: {output_file}")
        print("\n💡 Tips:")
        print("- Open the PowerPoint file in Microsoft PowerPoint or Google Slides")
        print("- You can customize colors, fonts, and layouts as needed")
        print("- Add images or animations to make it more engaging")
        
    except FileNotFoundError:
        print(f"❌ Error: Could not find {input_file}")
        print("Make sure the Markdown file exists in the current directory.")
    except ImportError:
        print("❌ Error: python-pptx library not found")
        print("Install it with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error: {e}")

if __name__ == "__main__":
    main() 